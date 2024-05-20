from pptx import Presentation
from pptx.util import Inches
from constants import CONSTANTS

# Create a presentation object
prs = Presentation()


# Function to set the background image for a slide
def set_background(slide, img_path):
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def create_slide(prs, layout, bg_img_path, title_text, content_text=None, bullet_points=None):
    slide_layout = prs.slide_layouts[layout]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, bg_img_path)

    title = slide.shapes.title
    title.text = title_text

    if content_text or bullet_points:
        content = slide.placeholders[1].text_frame
        content.text = content_text if content_text else ""

        if bullet_points:
            for point in bullet_points:
                p = content.add_paragraph()
                p.text = point


# Add a title slide
create_slide(prs,
             CONSTANTS['layout_type'][0],
             CONSTANTS['bg_img_path'],
             "About Myself",
             "A brief presentation about me")

# Add a slide with bullet points
create_slide(prs,
             CONSTANTS['layout_type'][1],
             CONSTANTS['bg_img_path'],
             "My personal information and skills",
             "Name: George",
             ["Age: 22", "Occupation: (hope to be an intern in Terralink)",
              "Hobbies: Coding, Adventure, Sports",
              "Skills: Python(even this presentation was written on Python!), "
              "Git/github, Cisco Packet Tracer, "
              "also i know English, and a bit German"])

# Add a slide with a list of skills
create_slide(prs,
             CONSTANTS['layout_type'][1],
             CONSTANTS['bg_img_path'],
             "My education",
             "Name: George",
             ["Moscow Aviation Institute (National Research University)",
              "Faculty: Control systems, computer science and electric power engineering",
              "Specialty: Computer Science and Engineering",
              "Education level: Bachelor's degree, "
              "“Automated management of business processes and finances”"])

create_slide(prs,
             CONSTANTS['layout_type'][0],
             CONSTANTS['bg_img_path'],
             "Internship in TerraLink",
             "Chosen Direction: System Analyst for the CDW team",
             ["Strong interest in data WH and system analysis",
              "Passionate about optimizing data systems for better performance",
              "Eager to learn from experienced professionals in the CDW team",
              "Excited to contribute to innovative projects"])

# Save the presentation
if __name__ == '__main__':
    prs.save('About_Myself.pptx')
